#!/usr/bin/env python3
"""
PowerPoint Modifier Script
Modifies the Global Capability Center presentation:
1. Changes footer from "Zinnov" to "SinVedha"
2. Applies complex design elements
3. Updates color scheme and styling

Requirements:
    pip install python-pptx pillow

Usage:
    python modify_powerpoint.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Configuration
INPUT_FILE = "Global capability Center_RAKSUL_Managed Services_Proposal.pptx"
OUTPUT_FILE = "Global_Capability_Center_RAKSUL_SinVedha_Complex_Design.pptx"

# Color Palette (Complex Design)
COLORS = {
    'deep_blue': RGBColor(26, 35, 126),      # #1a237e
    'teal': RGBColor(0, 105, 92),            # #00695c
    'coral': RGBColor(255, 111, 97),         # #ff6f61
    'light_gray': RGBColor(245, 245, 245),   # #f5f5f5
    'medium_gray': RGBColor(158, 158, 158),  # #9e9e9e
    'dark_gray': RGBColor(66, 66, 66),       # #424242
    'white': RGBColor(255, 255, 255),        # #ffffff
}

def modify_footer(prs):
    """Change footer from Zinnov to SinVedha"""
    print("Modifying footers...")
    
    # Modify slide master footers
    for master in prs.slide_master:
        for shape in master.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.lower()
                if 'zinnov' in text:
                    # Replace Zinnov with SinVedha
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if 'zinnov' in run.text.lower():
                                run.text = run.text.replace('Zinnov', 'SinVedha')
                                run.text = run.text.replace('zinnov', 'SinVedha')
                                run.text = run.text.replace('ZINNOV', 'SinVedha')
                                print(f"  ✓ Updated footer in master slide")
    
    # Modify individual slide footers
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.lower()
                if 'zinnov' in text:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if 'zinnov' in run.text.lower():
                                run.text = run.text.replace('Zinnov', 'SinVedha')
                                run.text = run.text.replace('zinnov', 'SinVedha')
                                run.text = run.text.replace('ZINNOV', 'SinVedha')
                                print(f"  ✓ Updated footer on slide {slide_num}")

def apply_complex_design(prs):
    """Apply complex design elements to slides"""
    print("\nApplying complex design elements...")
    
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\nProcessing slide {slide_num}...")
        
        # Add gradient background (simulated with colored rectangle)
        try:
            # Add a subtle background shape
            left = Inches(0)
            top = Inches(0)
            width = prs.slide_width
            height = prs.slide_height
            
            # Note: python-pptx doesn't support gradients directly
            # We'll add a semi-transparent shape instead
            bg_shape = slide.shapes.add_shape(
                1,  # Rectangle
                left, top, width, height
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = COLORS['light_gray']
            bg_shape.line.fill.background()
            
            # Send to back
            slide.shapes._spTree.remove(bg_shape._element)
            slide.shapes._spTree.insert(2, bg_shape._element)
            
            print(f"  ✓ Added background to slide {slide_num}")
        except Exception as e:
            print(f"  ⚠ Could not add background: {e}")
        
        # Update text formatting
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # Check if this is a title
                    if paragraph.level == 0:
                        for run in paragraph.runs:
                            run.font.size = Pt(44)
                            run.font.bold = True
                            run.font.color.rgb = COLORS['deep_blue']
                            print(f"  ✓ Styled title text")
                    else:
                        for run in paragraph.runs:
                            run.font.size = Pt(18)
                            run.font.color.rgb = COLORS['dark_gray']

def add_header_bar(slide, prs):
    """Add a colored header bar to slide"""
    try:
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = Inches(0.5)
        
        header = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS['deep_blue']
        header.line.fill.background()
        
        return header
    except Exception as e:
        print(f"  ⚠ Could not add header bar: {e}")
        return None

def enhance_slides(prs):
    """Add enhanced visual elements to slides"""
    print("\nEnhancing slides with visual elements...")
    
    for slide_num, slide in enumerate(prs.slides, 1):
        # Add header bar to content slides (skip title slide)
        if slide_num > 1:
            header = add_header_bar(slide, prs)
            if header:
                print(f"  ✓ Added header bar to slide {slide_num}")

def create_footer_text(prs):
    """Create consistent footer across all slides"""
    print("\nCreating consistent footer...")
    
    for slide_num, slide in enumerate(prs.slides, 1):
        try:
            # Add footer text box
            left = Inches(0.5)
            top = prs.slide_height - Inches(0.5)
            width = prs.slide_width - Inches(1)
            height = Inches(0.3)
            
            footer_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = footer_box.text_frame
            text_frame.clear()
            
            p = text_frame.paragraphs[0]
            p.text = f"SinVedha  |  Confidential  |  Page {slide_num}"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10)
            p.font.color.rgb = COLORS['medium_gray']
            
            print(f"  ✓ Added footer to slide {slide_num}")
        except Exception as e:
            print(f"  ⚠ Could not add footer to slide {slide_num}: {e}")

def main():
    """Main function to modify PowerPoint"""
    print("=" * 60)
    print("PowerPoint Design Modifier")
    print("=" * 60)
    
    # Check if input file exists
    if not os.path.exists(INPUT_FILE):
        print(f"\n❌ Error: File '{INPUT_FILE}' not found!")
        print("\nPlease ensure the PowerPoint file is in the same directory as this script.")
        return
    
    try:
        # Load presentation
        print(f"\n📂 Loading presentation: {INPUT_FILE}")
        prs = Presentation(INPUT_FILE)
        print(f"✓ Loaded {len(prs.slides)} slides")
        
        # Apply modifications
        modify_footer(prs)
        apply_complex_design(prs)
        enhance_slides(prs)
        create_footer_text(prs)
        
        # Save modified presentation
        print(f"\n💾 Saving modified presentation: {OUTPUT_FILE}")
        prs.save(OUTPUT_FILE)
        print(f"✓ Successfully saved!")
        
        print("\n" + "=" * 60)
        print("✅ COMPLETE!")
        print("=" * 60)
        print(f"\nYour modified presentation has been saved as:")
        print(f"  → {OUTPUT_FILE}")
        print("\nChanges made:")
        print("  ✓ Footer changed from 'Zinnov' to 'SinVedha'")
        print("  ✓ Applied complex color scheme")
        print("  ✓ Enhanced typography")
        print("  ✓ Added visual elements")
        print("\nPlease open the file in PowerPoint to review the changes.")
        
    except Exception as e:
        print(f"\n❌ Error: {e}")
        print("\nTroubleshooting:")
        print("  1. Ensure python-pptx is installed: pip install python-pptx")
        print("  2. Check that the PowerPoint file is not open in another program")
        print("  3. Verify you have write permissions in this directory")

if __name__ == "__main__":
    main()
